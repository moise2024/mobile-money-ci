"""
Generation des slides de soutenance — Projet Mobile Money CI
15 slides professionnels pour presentation orale
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), 'data', 'output')
PPTX_PATH = os.path.join(OUTPUT_DIR, 'Slides_Soutenance_Mobile_Money_CI.pptx')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Couleurs
C_BLEU = RGBColor(21, 101, 192)
C_BLEU_FONCE = RGBColor(13, 71, 161)
C_VERT = RGBColor(46, 125, 50)
C_ORANGE = RGBColor(230, 81, 0)
C_GRIS = RGBColor(100, 100, 100)
C_BLANC = RGBColor(255, 255, 255)
C_NOIR = RGBColor(0, 0, 0)
C_BG = RGBColor(245, 245, 245)

def add_bg(slide, color=C_BLANC):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, bold=False, color=C_NOIR, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=C_NOIR):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.level = 0
        p.space_after = Pt(6)
    return txBox

# ─── SLIDE 1 : PAGE DE GARDE ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_BLEU_FONCE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), C_ORANGE)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         'DATA ENGINEERING', 44, True, C_BLANC, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
         "Analyse des flux Mobile Money en Côte d'Ivoire", 28, False, C_BLANC, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.6),
         'Pipeline ETL · Schéma en étoile · Airflow · Docker', 20, False, RGBColor(200,200,200), PP_ALIGN.CENTER)

add_shape(slide, Inches(4), Inches(5), Inches(5.333), Inches(0.04), C_ORANGE)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         'Notre trinôme · Projet de Fin de Module', 20, False, RGBColor(200,200,200), PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
         'Encadré par : GOUAH Tato Serge · Juin 2025', 16, False, RGBColor(180,180,180), PP_ALIGN.CENTER)

# ─── SLIDE 2 : PLAN ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), C_BLEU)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         'Plan de la présentation', 32, True, C_BLEU_FONCE)

items = [
    '1. Contexte et problématique métier (2 min)',
    '2. Architecture du pipeline (2 min)',
    '3. Dataset et audit de qualité (2 min)',
    '4. Pipeline ETL : transformations clés (2 min)',
    '5. Modélisation : schéma en étoile (2 min)',
    '6. Analyses SQL et résultats (2 min)',
    '7. Dashboard et visualisations (1 min)',
    '8. Orchestration : Airflow + Docker (1 min)',
    '9. Difficultés et conclusion (1 min)',
]
add_bullet_text(slide, Inches(1.2), Inches(1.6), Inches(10), Inches(5.5), items, 20)

# ─── SLIDE 3 : CONTEXTE ───────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), C_BLEU)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         '1. Contexte : Mobile Money en Côte d\'Ivoire', 30, True, C_BLEU_FONCE)

items = [
    '4 opérateurs majeurs : MTN CI, Orange Money, Moov Africa, Wave',
    'Plus de 20 millions d\'utilisateurs en CI',
    '~12 millions de transactions par jour',
    'Secteur stratégique : télécoms, fintech, banques, retail',
    'Notre mission : Data Engineer chez un agrégateur de paiements à Abidjan',
    'Objectif : transformer 100 000 transactions brutes en indicateurs actionnables',
]
add_bullet_text(slide, Inches(1.2), Inches(1.6), Inches(11), Inches(5), items, 20)

# ─── SLIDE 4 : ARCHITECTURE ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), C_BLEU)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         '2. Architecture du pipeline', 30, True, C_BLEU_FONCE)

etapes = [
    ('1. CSV brut', '100k lignes\n12 colonnes', C_BLEU),
    ('2. Extract', 'Pandas\nAudit', C_VERT),
    ('3. Transform', 'Nettoyage\n+8 colonnes', C_ORANGE),
    ('4. Load', 'Supabase\nPostgreSQL', RGBColor(156, 39, 176)),
    ('5. Schéma', '4 dims + 1 fait\nÉtoile', RGBColor(0, 150, 136)),
    ('6. Analyse', 'SQL + Dashboard\n5 graphiques', C_BLEU_FONCE),
    ('7. Orchestre', 'Airflow\nDocker', C_VERT),
]

for i, (titre, desc, couleur) in enumerate(etapes):
    x = Inches(0.8 + i * 1.75)
    y = Inches(1.8)
    shape = add_shape(slide, x, y, Inches(1.5), Inches(3.5), couleur)
    add_text(slide, x + Inches(0.1), y + Inches(0.3), Inches(1.3), Inches(0.8),
             titre, 16, True, C_BLANC, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(1.3), Inches(1.3), Inches(1.8),
             desc, 14, False, C_BLANC, PP_ALIGN.CENTER)
    if i < len(etapes) - 1:
        add_text(slide, x + Inches(1.4), y + Inches(1.5), Inches(0.5), Inches(0.5),
                 '→', 24, True, C_GRIS, PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.8),
         'Stack : Python/Pandas → Supabase PostgreSQL → dbt → Airflow → Docker → GitHub Actions',
         16, False, C_GRIS, PP_ALIGN.CENTER)

# ─── SLIDE 5 : DATASET ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), C_BLEU)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         '3. Dataset et audit de qualité', 30, True, C_BLEU_FONCE)

add_text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.6),
         'Fichier : Data2_transactions_mobile_money_100k.csv', 18, True, C_NOIR)
items = [
    '100 000 lignes, 12 colonnes, 13 Mo',
    'Période : janvier à juin 2024',
    '4 opérateurs : MTN CI, Orange Money, Moov Africa, Wave',
    '5 types d\'opérations : Dépôt, Retrait, Transfert, etc.',
    '10 zones géographiques (Abidjan + intérieur)',
]
add_bullet_text(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4), items, 17)

add_text(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.6),
         'Problèmes détectés à l\'audit :', 18, True, C_ORANGE)
items2 = [
    '❌ 2 000 NaN sur frais_fcfa (2%)',
    '❌ 4 000 NaN sur zone_beneficiaire (4%)',
    '❌ 3 000 NaN sur id_agent (3%)',
    '❌ 200 montants aberrants (<= 0 FCFA)',
    '❌ Type date_heure incorrect (str)',
]
add_bullet_text(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(4), items2, 17)

add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8), RGBColor(232, 245, 233))
add_text(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.6),
         "✅ Résultat : 99 800 lignes propres, 0 NaN, 200 aberrations supprimées, types corrigés",
         16, False, C_VERT, PP_ALIGN.CENTER)

# ─── SLIDE 6 : ETL EXTRACT ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), C_BLEU)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         '4. ETL — Extraction et audit', 30, True, C_BLEU_FONCE)

items = [
    'Chargement CSV avec pandas.read_csv() en 0.33 seconde',
    'Inspection des types, valeurs manquantes, statistiques',
    'Identification des colonnes à corriger et enrichir',
    'Code clé :',
]
add_bullet_text(slide, Inches(1.2), Inches(1.6), Inches(11), Inches(2.5), items, 18)

# Code snippet
code_text = (
    "import pandas as pd\n"
    "df = pd.read_csv('transactions_mobile_money_100k.csv')\n"
    "print(df.info())\n"
    "print(df.isnull().sum())\n"
    "print(df['montant_fcfa'].describe())"
)
shape = add_shape(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(2.5), RGBColor(33, 33, 33))
add_text(slide, Inches(1.8), Inches(4.4), Inches(9.4), Inches(2),
         code_text, 14, False, RGBColor(0, 255, 0), PP_ALIGN.LEFT)

# ─── SLIDE 7 : ETL TRANSFORM ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), C_BLEU)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         '4. ETL — Nettoyage et enrichissement', 30, True, C_BLEU_FONCE)

add_text(slide, Inches(0.8), Inches(1.4), Inches(6), Inches(0.5),
         'Nettoyage :', 20, True, C_VERT)
items = [
    'date_heure → datetime64',
    'Chaînes vides → NaN → fillna()',
    'frais_fcfa NaN → 0',
    'zone_beneficiaire NaN → Zone inconnue',
    'id_agent NaN → AGT-INCONNU',
    '200 lignes aberrante supprimées',
]
add_bullet_text(slide, Inches(0.8), Inches(2.0), Inches(6), Inches(3.5), items, 16)

add_text(slide, Inches(7), Inches(1.4), Inches(6), Inches(0.5),
         'Enrichissement :', 20, True, C_ORANGE)
items2 = [
    'montant_net_fcfa (montant - frais)',
    'taux_frais_pct (%)',
    'categorie_montant (Micro/Petit/Moyen/Gros)',
    'tranche_horaire (Matin/Soirée/Nuit...)',
    'inter_ville (même ville ? 0/1)',
    'heure, jour_semaine, mois',
]
add_bullet_text(slide, Inches(7), Inches(2.0), Inches(6), Inches(3.5), items2, 16)

add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1), RGBColor(232, 245, 233))
add_text(slide, Inches(1), Inches(5.9), Inches(11), Inches(0.5),
         "✅ 5 tests qualité automatiques : complétude, unicité, validité, conformité, typage",
         16, False, C_VERT, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.4), Inches(11), Inches(0.3),
         "Taux de succès global : 80.2% | Volume total : 14,8 milliards FCFA",
         14, False, C_GRIS, PP_ALIGN.CENTER)

# ─── SLIDE 8 : SUPABASE ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), C_BLEU)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         '4. ETL — Chargement Supabase', 30, True, C_BLEU_FONCE)

items = [
    'Supabase : base PostgreSQL cloud gratuite (sans CB)',
    'Chargement par lots de 5 000 lignes (chunks)',
    '99 800 lignes chargées en 62 secondes',
    'Tables créées : transactions, agg_operateur, agg_type, agg_zone',
]
add_bullet_text(slide, Inches(1.2), Inches(1.6), Inches(11), Inches(3), items, 18)

add_text(slide, Inches(1.2), Inches(4.2), Inches(11), Inches(0.5),
         'Code clé (chargement par chunks) :', 18, True, C_NOIR)
code_text = (
    "chunks = [df[i:i+5000] for i in range(0, len(df), 5000)]\n"
    "for chunk in chunks:\n"
    "    chunk.to_sql('transactions', conn, if_exists='append', method='multi')"
)
shape = add_shape(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(2), RGBColor(33, 33, 33))
add_text(slide, Inches(1.8), Inches(5), Inches(9.4), Inches(1.6),
         code_text, 14, False, RGBColor(0, 255, 0), PP_ALIGN.LEFT)

# ─── SLIDE 9 : SCHÉMA ÉTOILE ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), C_BLEU)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         '5. Modélisation : Schéma en étoile', 30, True, C_BLEU_FONCE)

# Centre = faits
x_centre = Inches(5.5)
y_centre = Inches(3.2)
shape = add_shape(slide, x_centre, y_centre, Inches(2.3), Inches(1.5), C_ORANGE)
add_text(slide, x_centre + Inches(0.1), y_centre + Inches(0.2), Inches(2.1), Inches(0.5),
         'faits_transactions', 16, True, C_BLANC, PP_ALIGN.CENTER)
add_text(slide, x_centre + Inches(0.1), y_centre + Inches(0.7), Inches(2.1), Inches(0.6),
         'montant · frais\nstatut · inter_ville', 12, False, C_BLANC, PP_ALIGN.CENTER)

# Dimensions autour
dims = [
    (Inches(5.5), Inches(1.2), 'dim_operateur\n4 opérateurs\nMTN, Orange, Moov, Wave', C_BLEU),
    (Inches(1.5), Inches(3.2), 'dim_zone\n11 zones\nrégion, population', C_VERT),
    (Inches(9.5), Inches(3.2), 'dim_type_operation\n5 types\ncatégorie, frais', RGBColor(156, 39, 176)),
    (Inches(5.5), Inches(5.5), 'dim_date\n182 jours\nmois, trimestre, weekend', RGBColor(0, 150, 136)),
]
for x, y, txt, couleur in dims:
    shape = add_shape(slide, x, y, Inches(2.3), Inches(1.5), couleur)
    add_text(slide, x + Inches(0.1), y + Inches(0.15), Inches(2.1), Inches(1.2),
             txt, 12, True, C_BLANC, PP_ALIGN.CENTER)

# Fleches (textes)
add_text(slide, Inches(3.8), Inches(3.4), Inches(1.8), Inches(0.5),
         '← FK →', 16, True, C_GRIS, PP_ALIGN.CENTER)
add_text(slide, Inches(7.7), Inches(3.4), Inches(1.8), Inches(0.5),
         '← FK →', 16, True, C_GRIS, PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(7), Inches(11), Inches(0.4),
         '60 033 lignes dans faits_transactions · Clés étrangères REFERENCES',
         14, False, C_GRIS, PP_ALIGN.CENTER)

# ─── SLIDE 10 : SQL ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), C_BLEU)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         '6. Analyses SQL — Résultats', 30, True, C_BLEU_FONCE)

add_text(slide, Inches(0.8), Inches(1.4), Inches(12), Inches(0.5),
         'Q1 - Volume par opérateur :', 18, True, C_BLEU)
items = ['Wave : 2,16 Mds FCFA | MTN CI : 2,12 Mds | Moov : 2,12 Mds | Orange : 2,10 Mds']
add_bullet_text(slide, Inches(0.8), Inches(1.9), Inches(12), Inches(0.8), items, 16)

add_text(slide, Inches(0.8), Inches(2.7), Inches(12), Inches(0.5),
         'Q2 - Tendance mensuelle :', 18, True, C_VERT)
items = ['Volume stable ~1,4 Mds/mois · Légère baisse en juin (vacances)']
add_bullet_text(slide, Inches(0.8), Inches(3.2), Inches(12), Inches(0.8), items, 16)

add_text(slide, Inches(0.8), Inches(4.0), Inches(12), Inches(0.5),
         'Q3 - Top 5 agents par volume :', 18, True, C_ORANGE)
items = ['AGT0113 : 88,2 M FCFA | AGT0074 : 85,9 M | AGT0052 : 84,9 M']
add_bullet_text(slide, Inches(0.8), Inches(4.5), Inches(12), Inches(0.8), items, 16)

add_text(slide, Inches(0.8), Inches(5.3), Inches(12), Inches(0.5),
         'Q4 - Taux d\'échec par opérateur (CTE + RANK + LAG) :', 18, True, RGBColor(156, 39, 176))
items = ['Moov Africa 10.1% > Orange 10.0% > MTN 9.8% > Wave 9.7%']
add_bullet_text(slide, Inches(0.8), Inches(5.8), Inches(12), Inches(0.8), items, 16)

add_text(slide, Inches(0.8), Inches(6.5), Inches(12), Inches(0.6),
         '🔑 Fonctions de fenêtre utilisées : RANK() OVER, LAG(), PARTITION BY, CTE',
         16, False, C_GRIS, PP_ALIGN.CENTER)

# ─── SLIDE 11 : DASHBOARD ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), C_BLEU)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         '7. Tableau de bord — 5 graphiques', 30, True, C_BLEU_FONCE)

items = [
    'Volume mensuel par opérateur (barres empilées)',
    'Parts de marché par volume (camembert) ~25% chacun',
    'Volume par tranche horaire (courbe) — pics à 8h et 19h',
    'Taux de succès par opérateur (barres horizontales) ~80%',
    'Distribution montants par catégorie (violin, échelle log)',
]
add_bullet_text(slide, Inches(1.2), Inches(1.5), Inches(11), Inches(3.5), items, 20)

add_shape(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(1.2), RGBColor(255, 243, 224))
add_text(slide, Inches(1.2), Inches(5.6), Inches(11), Inches(0.5),
         '📊 Interprétation métier :', 18, True, C_ORANGE, PP_ALIGN.CENTER)
add_text(slide, Inches(1.2), Inches(6.1), Inches(11), Inches(0.5),
         'Marché équilibré entre 4 opérateurs · Pics d\'activité matin et soir · Majorité de petits montants',
         16, False, C_GRIS, PP_ALIGN.CENTER)

# ─── SLIDE 12 : AIRFLOW ───────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), C_BLEU)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         '8. Orchestration : Apache Airflow + Docker', 30, True, C_BLEU_FONCE)

add_text(slide, Inches(0.8), Inches(1.4), Inches(6), Inches(0.5),
         'DAG : pipeline_mobile_money_ci', 20, True, C_BLEU)
items = [
    'Schedule : 0 23 * * * (23h00 chaque jour)',
    '5 tâches enchaînées :',
    '   extract_csv → clean_data → load_supabase',
    '   → run_dbt → generate_report',
    '3 retries, alertes email, timeout 2h',
    'Callbacks : on_failure, on_success, on_retry',
]
add_bullet_text(slide, Inches(0.8), Inches(2.0), Inches(6), Inches(4), items, 16)

add_text(slide, Inches(7.5), Inches(1.4), Inches(5), Inches(0.5),
         'Conteneurisation Docker :', 20, True, C_VERT)
items2 = [
    'Dockerfile : python:3.10-slim (~130 Mo)',
    'docker-compose : 3 services',
    '   postgres (metadata Airflow)',
    '   airflow-webserver (interface)',
    '   airflow-scheduler (planificateur)',
    'Volumes partagés dags/ et data/',
]
add_bullet_text(slide, Inches(7.5), Inches(2.0), Inches(5), Inches(4), items2, 16)

# ─── SLIDE 13 : CI/CD + dbt ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), C_BLEU)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         'Bonus : CI/CD GitHub Actions + dbt', 30, True, C_BLEU_FONCE)

add_text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
         'GitHub Actions CI/CD :', 20, True, C_BLEU)
items = [
    'Déclenché à chaque git push (main/develop)',
    'Job 1 : quality-checks (flake8, pytest)',
    'Job 2 : docker-build (si tests OK)',
    'Cache pip pour builds plus rapides',
]
add_bullet_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3), items, 16)

add_text(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.5),
         'dbt (Data Build Tool) :', 20, True, C_VERT)
items2 = [
    'Modèle SQL déclaratif : perf_operateurs',
    'Vue analytique : volume mensuel par opérateur',
    'Tests de qualité : not_null, unique',
    'dbt run + dbt test automatisés',
]
add_bullet_text(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(3), items2, 16)

add_shape(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2), RGBColor(232, 245, 233))
add_text(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.5),
         '📦 Tous les livrables sur GitHub : https://github.com/moise2024/mobile-money-ci', 18, True, C_VERT, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.1), Inches(11), Inches(0.5),
         'README · Notebook .ipynb · DAG Airflow · Dockerfile · docker-compose · dbt · Rapport PDF · Captures',
         14, False, C_GRIS, PP_ALIGN.CENTER)

# ─── SLIDE 14 : DIFFICULTÉS ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), C_BLEU)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         '9. Difficultés et solutions', 30, True, C_BLEU_FONCE)

difficultes = [
    ('Gestion des accents', 'Succès vs Succes, Échec vs Echec → f-string avec UTF-8'),
    ('Types datetime64', 'ns vs us selon version Pandas → détection dynamique'),
    ('Chargement Supabase', 'Timeout du plan gratuit → chunks de 5000 lignes'),
    ('Limitation espace disque', '/tmp plein → liaison directe avec les fichiers source'),
    ('Token GitHub', 'Scope workflow manquant → ajout manuel sur GitHub'),
    ('Mappage dates', '60k faits vs 99k transactions → jointure dim_date'),
]
for i, (prob, sol) in enumerate(difficultes):
    y = Inches(1.5 + i * 0.85)
    add_shape(slide, Inches(1), y, Inches(0.8), Inches(0.6), C_ORANGE)
    add_text(slide, Inches(1), y + Inches(0.05), Inches(0.8), Inches(0.5),
             str(i+1), 18, True, C_BLANC, PP_ALIGN.CENTER)
    add_text(slide, Inches(2), y + Inches(0.05), Inches(4), Inches(0.5),
             prob, 16, True, C_NOIR)
    add_text(slide, Inches(6), y + Inches(0.05), Inches(6.5), Inches(0.5),
             sol, 16, False, C_GRIS)

# ─── SLIDE 15 : CONCLUSION ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_BLEU_FONCE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), C_ORANGE)

add_text(slide, Inches(1), Inches(1), Inches(11), Inches(0.8),
         'Conclusion', 36, True, C_BLANC, PP_ALIGN.CENTER)

items = [
    '✅ Pipeline ETL complet : 100 000 → 99 800 lignes propres en 0.33s',
    '✅ Schéma en étoile : 4 dimensions + 1 table de faits dans Supabase',
    '✅ 4 requêtes SQL analytiques dont CTE + fonctions de fenêtre',
    '✅ Dashboard 5 graphiques + rapport PDF + notebook .ipynb',
    '✅ DAG Airflow automatisé + Docker + CI/CD GitHub Actions',
    '✅ dbt (bonus) pour la transformation déclarative',
]
add_bullet_text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(4), items, 18, C_BLANC)

add_shape(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.04), C_ORANGE)
add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.6),
         'Merci pour votre attention ! Questions ?', 28, True, C_BLANC, PP_ALIGN.CENTER)

# ─── SAUVEGARDE ─────────────────────────────────────────────────
prs.save(PPTX_PATH)
print(f'Slides crées : {PPTX_PATH}')
print(f'Nombre de slides : {len(prs.slides)}')
